from pptx import Presentation
import re

from pathlib import Path
import requests
import shutil
import pandas as pd
import math

def find_file_name(url):
    return re.findall(r'(\w+.\w+)\?', url)[0]

def request_and_save_image(url, img_folder='images'):
    file_name = find_file_name(url)
    file_path_name = f'{img_folder}/{file_name}'
    if Path(file_path_name).exists():
        return file_path_name
    res = requests.get(url, stream = True)
    if res.status_code == 200:
        with open(file_path_name, 'wb') as f:
            shutil.copyfileobj(res.raw, f)
        return file_path_name


class SlidesBuilder:
    def __init__(self, dataframe):
        self.df = dataframe
        self.template_path = 'template/template.pptx'
    
    def buildSlides(self):
        prs = Presentation(self.template_path)
        ITEMS_PER_SLIDE = 8

        for collection_name, data in self.df.groupby('collection'):
            grouped_df = pd.DataFrame(data)
            names = grouped_df['name'].to_list()
            prices = grouped_df['price'].to_list()
            image_links = grouped_df['img'].to_list()
            total_item_count = len(grouped_df)
            number_collection_slides = math.ceil(total_item_count / ITEMS_PER_SLIDE)

            #add collection slide [section slide]
            collection_layout = prs.slide_layouts[0]
            coll_slide = prs.slides.add_slide(collection_layout)
            for shape in coll_slide.placeholders:
                if shape.is_placeholder:
                    shape.text = collection_name

            image_shape_list = list()  # list of image of each slide len(image_shape_list) == ITEMS_PER_SLIDE
            name_shape_list = list()
            price_shape_list = list()

            # individual catalogue layout 
            for slide in range(number_collection_slides):
                # for each slide
                catalogue_layout= prs.slide_layouts[1]
                cat_slide = prs.slides.add_slide(catalogue_layout)
                is_name = True
                for shape in cat_slide.placeholders:
                    if shape.is_placeholder:
                        placeholder_type = shape.placeholder_format.type
                        if placeholder_type == 18:
                            image_shape_list.append(shape)
                        elif is_name:
                            name_shape_list.append(shape)
                            is_name =  False
                        else:
                            price_shape_list.append(shape)
                            is_name = True

            for i, name in enumerate(names):
                price = prices[i]
                image_link = image_links[i]
                img_path = request_and_save_image(url=image_link)
                print(f"PROCESSING: {name}, {price}, {img_path}")
                image_shape_list[i].insert_picture(img_path)
                name_shape_list[i].text = name
                price_shape_list[i].text = price

        prs.save('products.pptx')