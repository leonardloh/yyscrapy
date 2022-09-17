from pptx import Presentation
from pptx.util import Inches

from pathlib import Path
import requests
import shutil


# image_folder_path = Path('images')
# url = df.iloc[0, 3]
# res = requests.get(url, stream = True)
# if res.status_code == 200:
#     with open(f'{image_folder_path}/{df.iloc[0, 1]}.jpg', 'wb') as f:
#         shutil.copyfileobj(res.raw, f)

def request_and_save_image(img_name, url, img_folder='images'):
    res = requests.get(url, stream = True)
    file_path_name = f'{img_folder}/{img_name}.jpg'
    if res.status_code == 200:
        with open(file_path_name, 'wb') as f:
            shutil.copyfileobj(res.raw, f)
        return file_path_name


class SlidesBuilder:
    def __init__(self, dataframe):
        self.df = dataframe
    
    def buildSlides(self):
        

        collection = self.df.iloc[0, 0]
        name = self.df.iloc[0, 1]
        price = self.df.iloc[0, 2]
        img_link = self.df.iloc[0, 3]
        img_path = request_and_save_image(name, img_link)

        prs = Presentation('template.pptx')
        collection_layout = prs.slide_layouts[0]
        catalogue_layout= prs.slide_layouts[1]

        coll_slide = prs.slides.add_slide(collection_layout)
        cat_slide = prs.slides.add_slide(catalogue_layout)

        for shape in coll_slide.placeholders:
            if shape.is_placeholder:
                #  print(f'{shape.placeholder_format.idx, shape.name, shape.placeholder_format.type}')
                shape.text = collection

        is_name = True
        for shape in cat_slide.placeholders:
            if shape.is_placeholder:
                # print(f'{shape.placeholder_format.idx, shape.name, shape.placeholder_format.type}')
                placeholder_type = shape.placeholder_format.type
                # if image
                if placeholder_type == 18:
                    shape.insert_picture(img_path)
                elif is_name:
                    shape.text = name
                    is_name =  False
                else:
                    shape.text = price
                    is_name = True


        # cat_slide.placeholders[10].insert_picture(img_path)
        # cat_slide.placeholders[11].text = name
        # cat_slide.placeholders[12].text = price



        prs.save('test.pptx')